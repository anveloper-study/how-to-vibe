"""
HTML 슬라이드 → 편집 가능한 PPTX 변환 스크립트
- data-object-type="textbox" → PPTX 텍스트박스
- data-object-type="shape" → PPTX 도형 (사각형)
- 절대 좌표(px) → PPTX EMU 단위 변환
"""

import os
import re
import glob
from bs4 import BeautifulSoup, NavigableString
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# 기본 폰트
DEFAULT_FONT = 'Noto Sans KR'

# 1280x720 HTML → 표준 16:9 슬라이드 (10x5.625 inches)
SLIDE_W_PX = 1280
SLIDE_H_PX = 720
SLIDE_W_EMU = Inches(10)
SLIDE_H_EMU = Inches(5.625)

def px_to_emu_x(px):
    return int(px / SLIDE_W_PX * SLIDE_W_EMU)

def px_to_emu_y(px):
    return int(px / SLIDE_H_PX * SLIDE_H_EMU)

def px_to_emu_w(px):
    return int(px / SLIDE_W_PX * SLIDE_W_EMU)

def px_to_emu_h(px):
    return int(px / SLIDE_H_PX * SLIDE_H_EMU)

def parse_css_value(style_str, prop):
    """CSS 속성값 추출"""
    pattern = rf'{prop}\s*:\s*([^;]+)'
    m = re.search(pattern, style_str)
    if m:
        return m.group(1).strip()
    return None

def parse_px(val):
    """'123px' → 123"""
    if val is None:
        return 0
    m = re.search(r'(-?\d+(?:\.\d+)?)\s*px', str(val))
    return float(m.group(1)) if m else 0

def parse_color(color_str):
    """CSS color → RGBColor"""
    if not color_str:
        return None
    color_str = color_str.strip()

    # hex 6자리
    m = re.match(r'#([0-9a-fA-F]{6})', color_str)
    if m:
        h = m.group(1)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    # hex 3자리
    m = re.match(r'#([0-9a-fA-F]{3})$', color_str)
    if m:
        h = m.group(1)
        return RGBColor(int(h[0]*2, 16), int(h[1]*2, 16), int(h[2]*2, 16))

    # rgb/rgba
    m = re.match(r'rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)', color_str)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))

    return None

def parse_opacity(color_str):
    """rgba에서 opacity 추출"""
    if not color_str:
        return 1.0
    m = re.match(r'rgba\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*,\s*([\d.]+)\s*\)', color_str)
    if m:
        return float(m.group(1))
    return 1.0

def parse_font_size(val):
    """'16px' → Pt(12) (px * 0.75 = pt)"""
    if not val:
        return Pt(10.5)
    m = re.search(r'(\d+(?:\.\d+)?)\s*px', str(val))
    if m:
        return Pt(float(m.group(1)) * 0.75)
    return Pt(10.5)

def parse_border_radius(style_str):
    """border-radius 파싱"""
    val = parse_css_value(style_str, 'border-radius')
    if val:
        return parse_px(val)
    return 0

def set_shape_transparency(shape, opacity):
    """도형의 fill에 투명도(alpha) 설정"""
    if opacity >= 1.0:
        return
    alpha_val = int(opacity * 100000)
    # spPr → solidFill → srgbClr 접근
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        return
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        for old in srgb.findall(qn('a:alpha')):
            srgb.remove(old)
        alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
        alpha_elem.set('val', str(alpha_val))

def set_line_transparency(shape, opacity):
    """테두리 투명도 설정"""
    if opacity >= 1.0:
        return
    alpha_val = int(opacity * 100000)
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        return
    solid = ln.find(qn('a:solidFill'))
    if solid is None:
        return
    srgb = solid.find(qn('a:srgbClr'))
    if srgb is not None:
        for old in srgb.findall(qn('a:alpha')):
            srgb.remove(old)
        alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
        alpha_elem.set('val', str(alpha_val))

def get_element_bounds(el):
    """요소의 위치/크기 추출"""
    style = el.get('style', '')
    left = parse_px(parse_css_value(style, 'left'))
    top = parse_px(parse_css_value(style, 'top'))
    width = parse_px(parse_css_value(style, 'width'))
    height = parse_px(parse_css_value(style, 'height'))
    return left, top, width, height

def extract_text_runs(element):
    """HTML 요소에서 텍스트 런(스타일별 텍스트 조각) 추출"""
    runs = []

    def walk(node, inherited_style=""):
        if isinstance(node, NavigableString):
            text = str(node)
            text = text.replace('\n', ' ').strip()
            if text:
                runs.append({'text': text, 'style': inherited_style})
            return

        if node.name == 'br':
            runs.append({'text': '\n', 'style': '', 'is_break': True})
            return

        style = node.get('style', '')
        combined = inherited_style + '; ' + style if inherited_style else style

        has_children = False
        for child in node.children:
            has_children = True
            walk(child, combined)

        if not has_children and node.string:
            text = node.string.replace('\n', ' ').strip()
            if text:
                runs.append({'text': text, 'style': combined})

    p_tags = element.find_all('p', recursive=True)
    if p_tags:
        for i, p in enumerate(p_tags):
            if i > 0:
                runs.append({'text': '\n', 'style': '', 'is_break': True})
            p_style = p.get('style', '')
            for child in p.children:
                walk(child, p_style)
    else:
        for child in element.children:
            walk(child, '')

    return runs

def apply_run_style(run_obj, style_str):
    """pptx run에 CSS 스타일 적용"""
    font = run_obj.font

    # 기본 폰트
    font.name = DEFAULT_FONT

    # color
    color = parse_css_value(style_str, 'color')
    if color:
        rgb = parse_color(color)
        if rgb:
            font.color.rgb = rgb
    else:
        # 기본 흰색 (어두운 배경)
        font.color.rgb = RGBColor(255, 255, 255)

    # font-size
    fs = parse_css_value(style_str, 'font-size')
    if fs:
        font.size = parse_font_size(fs)

    # font-weight
    fw = parse_css_value(style_str, 'font-weight')
    if fw and (fw in ('500', '700', '800', '900', 'bold')):
        font.bold = True

    # letter-spacing (PPTX에서는 spacing으로)
    ls = parse_css_value(style_str, 'letter-spacing')
    if ls:
        ls_px = parse_px(ls)
        if ls_px != 0:
            font.spacing = Pt(ls_px * 0.75)

    # font-family
    ff = parse_css_value(style_str, 'font-family')
    if ff:
        if 'monospace' in ff.lower() or 'courier' in ff.lower():
            font.name = 'Courier New'

def get_text_align(style_str):
    """text-align → PP_ALIGN"""
    ta = parse_css_value(style_str, 'text-align')
    if ta == 'center':
        return PP_ALIGN.CENTER
    elif ta == 'right':
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT

def add_textbox(slide, el):
    """HTML textbox → PPTX 텍스트박스"""
    left, top, width, height = get_element_bounds(el)
    if width == 0 or height == 0:
        return

    txBox = slide.shapes.add_textbox(
        px_to_emu_x(left), px_to_emu_y(top),
        px_to_emu_w(width), px_to_emu_h(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    runs_data = extract_text_runs(el)
    if not runs_data:
        return

    first_p = el.find('p')
    p_style = first_p.get('style', '') if first_p else ''

    para = tf.paragraphs[0]
    para.alignment = get_text_align(p_style)
    para.space_before = Pt(0)
    para.space_after = Pt(0)

    for rd in runs_data:
        if rd.get('is_break'):
            para = tf.add_paragraph()
            para.alignment = get_text_align(p_style)
            para.space_before = Pt(0)
            para.space_after = Pt(0)
            continue

        text = rd['text']
        if not text.strip():
            continue

        run = para.add_run()
        run.text = text
        apply_run_style(run, rd['style'])

def parse_border_style(style_str):
    """border 관련 스타일 파싱 → (color, width, opacity)"""
    results = []

    # border-top 우선
    for prop in ['border-top', 'border-left', 'border']:
        val = parse_css_value(style_str, prop)
        if val and 'solid' in val:
            parts = val.split('solid')
            width_part = parts[0].strip()
            color_part = parts[1].strip() if len(parts) > 1 else ''
            bw = parse_px(width_part) if width_part else 1
            bc = parse_color(color_part)
            bo = parse_opacity(color_part) if 'rgba' in color_part else 1.0
            if bc:
                results.append((prop, bc, bw, bo))

    return results

def add_shape(slide, el):
    """HTML shape → PPTX 사각형 도형"""
    left, top, width, height = get_element_bounds(el)
    if width <= 6 or height <= 6:
        return  # 장식 점 스킵

    style = el.get('style', '')
    bg_color_str = parse_css_value(style, 'background-color')
    opacity = 1.0
    rgb = None

    if bg_color_str:
        opacity = parse_opacity(bg_color_str)
        rgb = parse_color(bg_color_str)

    # 선(height <= 4 or width <= 4)은 라인으로 처리
    if height <= 4 and width > 10:
        # 수평선
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            px_to_emu_x(left), px_to_emu_y(top),
            px_to_emu_w(width), px_to_emu_h(height)
        )
        if rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb
            if opacity < 1.0:
                set_shape_transparency(shape, opacity)
        else:
            shape.fill.background()
        shape.line.fill.background()
        return

    if width <= 4 and height > 10:
        # 수직선
        shape = slide.shapes.add_shape(
            1,
            px_to_emu_x(left), px_to_emu_y(top),
            px_to_emu_w(width), px_to_emu_h(height)
        )
        if rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb
            if opacity < 1.0:
                set_shape_transparency(shape, opacity)
        else:
            shape.fill.background()
        shape.line.fill.background()
        return

    from pptx.enum.shapes import MSO_SHAPE
    radius = parse_border_radius(style)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius > 8 else MSO_SHAPE.RECTANGLE,
        px_to_emu_x(left), px_to_emu_y(top),
        px_to_emu_w(width), px_to_emu_h(height)
    )

    # 배경색 + 투명도
    if rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
        if opacity < 1.0:
            set_shape_transparency(shape, opacity)
    else:
        shape.fill.background()

    # 테두리 처리
    borders = parse_border_style(style)
    if borders:
        # 가장 두드러진 테두리 사용
        _, bc, bw, bo = borders[0]
        shape.line.color.rgb = bc
        shape.line.width = Pt(bw)
        if bo < 1.0:
            set_line_transparency(shape, bo)
    else:
        shape.line.fill.background()

    # box-shadow가 있으면 글로우 효과 (PPTX에서는 shadow로)
    if 'box-shadow' in style and 'FE2A80' in style.upper().replace(' ', ''):
        shape.shadow.inherit = False

def process_slide(prs, html_path):
    """HTML 파일 1개 → PPTX 슬라이드 1개"""
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'lxml')

    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # 슬라이드 배경을 검정으로
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 10)

    # data-object 요소들 수집
    objects = soup.find_all(attrs={"data-object": "true"})

    # z-index로 정렬
    def get_z(el):
        style = el.get('style', '')
        z = parse_css_value(style, 'z-index')
        return int(z) if z and z.isdigit() else 0

    objects.sort(key=get_z)

    for obj in objects:
        obj_type = obj.get('data-object-type', '')
        style = obj.get('style', '')

        # blur 글로우 효과는 스킵
        if 'blur' in style:
            continue

        if obj_type == 'shape':
            add_shape(slide, obj)
        elif obj_type == 'textbox':
            add_textbox(slide, obj)

def main():
    slide_dir = os.path.join(os.path.dirname(__file__), '..', 'docs', 'slide')
    output_path = os.path.join(os.path.dirname(__file__), '..', 'docs', 'vibe-coding-slides.pptx')

    html_files = sorted(glob.glob(os.path.join(slide_dir, '*.html')))
    print(f"발견된 슬라이드: {len(html_files)}개")

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU

    for html_file in html_files:
        name = os.path.basename(html_file)
        print(f"  변환 중: {name}")
        process_slide(prs, html_file)

    prs.save(output_path)
    print(f"\n완료: {output_path}")

if __name__ == '__main__':
    main()
